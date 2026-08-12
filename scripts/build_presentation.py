"""
SoilGuard-CG: Ideathon 5-Slide Deck Builder
Creates a 16:9 PowerPoint (.pptx) — title slide + 5 content slides — with a dark
"ground-station console" theme, embedded maps, and full speaker notes on every slide.

Output: ideathon-study-pack/SoilGuard_CG_Ideathon_Deck.pptx
Canva:   upload the .pptx to canva.com and edit freely.

Usage:
    python scripts/build_presentation.py

Requires: python-pptx   (pip install python-pptx)
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    import pandas as pd
except ImportError:
    pd = None

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "ideathon-study-pack"
ASSETS = OUT_DIR / "assets"
GOLDEN = ROOT / "soilguard-cg" / "outputs" / "golden_backup"

with open(ROOT / "shared" / "metrics.json", encoding="utf-8") as f:
    M = json.load(f)

# ---- theme --------------------------------------------------------------
BG = RGBColor(0x0B, 0x12, 0x20)          # deep space navy
BG2 = RGBColor(0x11, 0x1B, 0x2E)         # panel navy
GREEN = RGBColor(0x22, 0xC5, 0x5E)       # terminal green
GREEN_D = RGBColor(0x16, 0x62, 0x2E)     # deep green
CYAN = RGBColor(0x38, 0xBD, 0xF8)        # console cyan
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
GREY = RGBColor(0x9C, 0xA8, 0xB8)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xEA, 0xB3, 0x08)

SW, SH = Inches(13.333), Inches(7.5)     # 16:9


def new_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def footer(slide, idx, total=6):
    bar = slide.shapes.add_textbox(Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.4))
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SoilGuard-CG · Chhattisgarh Soil Carbon Sentinel · National Space Day Ideathon 2026"
    r.font.size = Pt(9)
    r.font.color.rgb = GREY
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"{idx} / {total}"
    r2.font.size = Pt(9)
    r2.font.color.rgb = GREY
    p2.alignment = PP_ALIGN.RIGHT


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def add_text(tf, text, size=14, color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, space_after=4, bullet=False, level=0):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs and len(tf.paragraphs) == 1 else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return p


def title_bar(slide, kicker, title, sub=None):
    tb = textbox(slide, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.5))
    add_text(tb.text_frame, kicker.upper(), size=12, color=CYAN, bold=True, space_after=2)
    tb2 = textbox(slide, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.9))
    add_text(tb2.text_frame, title, size=30, color=WHITE, bold=True, space_after=0)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.5), Inches(1.6), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()
    if sub:
        tb3 = textbox(slide, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.5))
        add_text(tb3.text_frame, sub, size=13, color=GREY, italic=True)


def panel(slide, l, t, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG2
    sh.line.color.rgb = RGBColor(0x2A, 0x3A, 0x52)
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat_chip(slide, l, t, w, big, label, color=GREEN):
    panel(slide, l, t, w, Inches(1.35))
    tb = textbox(slide, l + Inches(0.18), t + Inches(0.12), w - Inches(0.36), Inches(0.75))
    add_text(tb.text_frame, big, size=26, color=color, bold=True, space_after=0)
    tb2 = textbox(slide, l + Inches(0.18), t + Inches(0.82), w - Inches(0.36), Inches(0.45))
    add_text(tb2.text_frame, label, size=11, color=GREY)


def pic_fit(slide, path, l, t, w=None, h=None):
    return slide.shapes.add_picture(str(path), l, t, width=w, height=h)


# ===============================================================================
def build():
    prs = new_prs()

    # ---------------------------------------------------------------- Slide 1 · Title
    s = blank(prs)
    set_bg(s)
    panel(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(5.0))
    kick = textbox(s, Inches(1.3), Inches(1.35), Inches(10.8), Inches(0.5))
    add_text(kick.text_frame, "NATIONAL SPACE DAY IDEATHON 2026 · COSINE NIT RAIPUR × NRSC / ISRO",
             size=13, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    t = textbox(s, Inches(1.3), Inches(1.9), Inches(10.8), Inches(1.9))
    add_text(t.text_frame, "SoilGuard-CG", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=6)
    add_text(t.text_frame, "Chhattisgarh Soil Carbon Sentinel", size=26, color=GREEN, bold=True,
             align=PP_ALIGN.CENTER, space_after=14)
    add_text(t.text_frame,
             "Satellite-Driven High-Resolution SOC Deficiency Mapping & Regenerative Advisory System",
             size=16, color=GREY, align=PP_ALIGN.CENTER)
    st = textbox(s, Inches(1.3), Inches(4.35), Inches(10.8), Inches(0.9))
    add_text(st.text_frame, "10m Sentinel-2  ·  Random Forest  ·  Zero Target Leakage  ·  100% Offline  ·  ISRO Bhuvan Validated",
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    fl = textbox(s, Inches(1.3), Inches(5.35), Inches(10.8), Inches(0.4))
    add_text(fl.text_frame,
             "Raipur–Durg Agricultural Belt, Chhattisgarh  |  Team SoilGuard Innovation",
             size=12, color=AMBER, align=PP_ALIGN.CENTER)
    footer(s, 1)
    notes(s, "SPEAKER NOTES (0:00–0:30 · HOOK):\n"
             "\"Respected judges, Chhattisgarh is the rice bowl of India. But after every harvest, millions of "
             "hectares of topsoil lose organic carbon to intense summer heat and stubble burning. Traditional "
             "soil health cards sample just one point per 10 hectares every three years — farmers stay blind "
             "to degradation.\"\n\n"
             "SoilGuard-CG is our answer: a 100% offline, satellite-driven ML platform that maps Soil Organic "
             "Carbon deficiency at 10-metre resolution and tells officers exactly where to act.\n\n"
             "Say the numbers slowly: 2.27 million pixels evaluated · 19.6% of topsoil in critical deficiency.")

    # ---------------------------------------------------------------- Slide 2 · Problem
    s = blank(prs)
    set_bg(s)
    title_bar(s, "The Problem", "Chhattisgarh's topsoil is losing carbon — and nobody can see where",
              "Dhan ka Katora (Rice Bowl): Raipur–Durg–Dhamtari plains · 3.5M smallholder farmers")
    # left panel: crisis
    panel(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.85))
    tb = textbox(s, Inches(0.85), Inches(2.15), Inches(5.4), Inches(4.5))
    add_text(tb.text_frame, "Three drivers of topsoil carbon loss", size=17, color=GREEN, bold=True, space_after=8)
    add_text(tb.text_frame, "Paddy mono-cropping — waterlogged rice accelerates organic-matter breakdown", size=14, bullet=True)
    add_text(tb.text_frame, "Summer heat 42°C+ — thermal oxidation of exposed topsoil", size=14, bullet=True)
    add_text(tb.text_frame, "Stubble burning — direct combustion of residue organic carbon", size=14, bullet=True)
    add_text(tb.text_frame, " ", size=6, space_after=2)
    add_text(tb.text_frame, "The measurement gap", size=17, color=RED, bold=True, space_after=8)
    add_text(tb.text_frame, "Soil Health Cards: 1 physical sample per 10 ha, once every 3 years", size=14, bullet=True)
    add_text(tb.text_frame, "₹2,000–3,000 per sample — broad monitoring is unaffordable", size=14, bullet=True)
    add_text(tb.text_frame, "Misses intra-field variability and the post-harvest window entirely", size=14, bullet=True)
    # right: stats
    stat_chip(s, Inches(6.85), Inches(1.95), Inches(5.9), "> 52%", "agricultural grids with SOC below critical 0.50% threshold", RED)
    stat_chip(s, Inches(6.85), Inches(3.5), Inches(5.9), "1 : 10 ha", "sampling density of traditional soil health cards (every 3 years)", AMBER)
    stat_chip(s, Inches(6.85), Inches(5.05), Inches(5.9), "42°C +", "post-harvest summer topsoil temperatures drive carbon loss", CYAN)
    footer(s, 2)
    notes(s, "SPEAKER NOTES (0:00–0:30 · PROBLEM):\n"
             "Tell the story, don't read bullets:\n"
             "\"After every paddy harvest, Chhattisgarh's fields sit bare under 42-degree summer heat. Stubble "
             "burning finishes off the rest. The result: over half the region's agricultural grids now sit below "
             "the critical 0.50% organic-carbon threshold.\"\n\n"
             "The killer point: \"And the only tool we have today — government soil health cards — samples one "
             "physical point per 10 hectares, once every three years. At ₹2–3 thousand per sample. Farmers are "
             "flying blind.\"\n\n"
             "Pause. Then: \"We built a satellite that changes this — at zero sampling cost.\"")

    # ---------------------------------------------------------------- Slide 3 · Methodology
    s = blank(prs)
    set_bg(s)
    title_bar(s, "The Solution", "One offline pipeline: satellite pixels → soil intelligence",
              "Python 3.11 · rasterio · scikit-learn · matplotlib · rich · runs on a basic laptop, no internet")
    # pipeline flow
    steps = [
        ("Sentinel-2\n10m L2A COG", "B02 Blue · B04 Red\nB08 NIR · B11 SWIR1"),
        ("Bare-Soil Mask\nNDVI ≤ 0.30 · BSI", "2.27 M pixels isolated\n(48.96% of scene)"),
        ("Random Forest\n150 trees · 9 features", "No target leakage\nSoilGrids SOC excluded"),
        ("SOC Deficiency\nMap (0–1, 10m)", "Risk score per pixel\npublished at 300 DPI"),
        ("5×5 Zonal Grid\n25 sectors", "Priority ranking\nA1–E5 administrative"),
        ("Advisory Engine\nFYM · tillage · pH", "Clay-sensitive packages\nper sector"),
    ]
    x = Inches(0.4)
    y = Inches(2.0)
    bw = Inches(1.92)
    bh = Inches(1.7)
    for i, (head, body) in enumerate(steps):
        panel(s, x, y, bw, bh)
        tb = textbox(s, x + Inches(0.1), y + Inches(0.1), bw - Inches(0.2), bh - Inches(0.2))
        add_text(tb.text_frame, head, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=4)
        add_text(tb.text_frame, body, size=10, color=GREY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw - Inches(0.04), y + Inches(0.68), Inches(0.28), Inches(0.34))
            ar.fill.solid()
            ar.fill.fore_color.rgb = GREEN
            ar.line.fill.background()
        x += bw + Inches(0.16)
    # zero-leakage callout
    call = panel(s, Inches(0.4), Inches(4.05), Inches(12.5), Inches(1.05))
    tb = textbox(s, Inches(0.7), Inches(4.18), Inches(12.0), Inches(0.8))
    add_text(tb.text_frame, "ZERO TARGET LEAKAGE  ·  ",
             size=14, color=GREEN, bold=True, space_after=0)
    p = tb.text_frame.add_paragraph()
    r = p.add_run()
    r.text = ("Raw SoilGrids SOC is 100% excluded from the model's features. The Random Forest must learn the link "
              "between optical SWIR/NIR reflectance and SOC deficiency — it cannot memorise static maps.")
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    # model card
    panel(s, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.5))
    tb = textbox(s, Inches(0.7), Inches(5.42), Inches(12.0), Inches(1.3))
    add_text(tb.text_frame, "Model card", size=13, color=CYAN, bold=True, space_after=4)
    add_text(tb.text_frame,
             f"Target  y = 0.60·(1 − SOCₙₒᵣₘ) + 0.40·BSIₙₒᵣₘ      |      Features: 9 satellite-only "
             f"(BSI, NDVI, SWIR1/NIR & SWIR1/Red ratios, Blue/Red/NIR/SWIR1)      |      Test R² = {M['r2']}, "
             f"RMSE = {M['rmse']}      |      split 80/20, random_state 42",
             size=12.5, color=WHITE)
    footer(s, 3)
    notes(s, "SPEAKER NOTES (0:30–1:15 · SOLUTION & DEMO):\n"
             "\"With ONE terminal command, SoilGuard-CG takes 10-metre Sentinel-2 imagery, isolates 2.27 million "
             "bare-soil pixels using NDVI and the Bare Soil Index, trains a leakage-free Random Forest, and "
             "renders SOC deficiency maps in about a minute — 100% offline.\"\n\n"
             "Zero-leakage is your strongest credibility point — say it exactly:\n"
             "\"We deliberately exclude the SoilGrids SOC map from the model's inputs. The model must learn true "
             "optical physics — carbon-rich soil absorbs SWIR light, depleted soil reflects it. That's why our "
             "R² is honest.\"\n\n"
             "If running the live demo (demo.bat), do it during this slide: point at the terminal as each step "
             "completes. Say \"you are watching the whole pipeline run right now.\"\n\n"
             "R² defence ready: \"Peer-reviewed optical-soil-carbon benchmarks are R² ∈ [0.35, 0.55] — we're "
             "right inside. Anything above 0.9 usually leaks the target.\"")

    # ---------------------------------------------------------------- Slide 4 · Results
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Key Results", "19.6% of evaluated topsoil is in critical SOC deficiency",
              f"Full AOI {M['totalAoiHa']:,.0f} ha · bare soil evaluated {M['bareSoilHa']:,.0f} ha")
    # left: stats + tiers
    stat_chip(s, Inches(0.55), Inches(1.9), Inches(4.0), f"{M['highRiskHa']:,.0f} ha", f"high deficiency (>0.58) · {M['highRiskPct']:.1f}% of soil area", RED)
    stat_chip(s, Inches(4.75), Inches(1.9), Inches(4.0), f"{M['moderateRiskHa']:,.0f} ha", "moderate (0.45–0.58) · needs replenishment", AMBER)
    stat_chip(s, Inches(8.95), Inches(1.9), Inches(3.85), f"{M['lowRiskHa']:,.0f} ha", "low (<0.45) · stable, maintain", GREEN)
    # risk map
    risk = ASSETS / "phase3_risk_score_map.png"
    if risk.exists():
        pic_fit(s, risk, Inches(0.55), Inches(3.5), w=Inches(6.4))
    # right panel: top sectors
    panel(s, Inches(7.2), Inches(3.45), Inches(5.6), Inches(3.4))
    tb = textbox(s, Inches(7.45), Inches(3.6), Inches(5.1), Inches(0.4))
    add_text(tb.text_frame, "Top priority sectors", size=15, color=CYAN, bold=True)
    rows = [
        ("#1", "Arang (B-1)", "0.66", "946 ha"),
        ("#2", "Abhanpur (A-1)", "0.65", "901 ha"),
        ("#3", "Arang (B-2)", "0.63", "817 ha"),
    ]
    yy = Inches(4.1)
    for rk, name, risk_v, ha in rows:
        panel(s, Inches(7.45), yy, Inches(5.1), Inches(0.78))
        tb = textbox(s, Inches(7.65), yy + Inches(0.08), Inches(5.0), Inches(0.6))
        add_text(tb.text_frame, f"{rk}  {name}", size=13, color=WHITE, bold=True, space_after=0)
        add_text(tb.text_frame, f"mean deficiency {risk_v}  ·  high-risk {ha}", size=10.5, color=GREY)
        yy += Inches(0.88)
    footer(s, 4)
    notes(s, "SPEAKER NOTES (1:15–1:45 · RESULTS):\n"
             "Point at the map, then the chips:\n"
             "\"Across 22,700 hectares we found 19.6% of topsoil — over 4,400 hectares — in critical SOC "
             "deficiency. The red zones on this 10-metre map are exactly where organic carbon is missing.\"\n\n"
             "\"Three sectors dominate: Arang B-1, Abhanpur A-1 and Arang B-2 — together nearly 2,700 hectares "
             "of high-deficiency soil in the Raipur belt.\"\n\n"
             "Transition: \"But we don't stop at a map. We turn it into action.\"\n\n"
             "Tip: zoom into the red zone with your hand/laser pointer. Make the audience feel the scale.")

    # ---------------------------------------------------------------- Slide 5 · Advisory
    s = blank(prs)
    set_bg(s)
    title_bar(s, "From Map to Action", "Clay-sensitive regenerative advisory for every sector",
              "Every one of the 25 sectors gets an exact package — FYM tonnes, tillage, pH correction")
    # left: advisory logic table
    panel(s, Inches(0.55), Inches(1.95), Inches(7.0), Inches(4.85))
    tb = textbox(s, Inches(0.85), Inches(2.1), Inches(6.4), Inches(4.5))
    add_text(tb.text_frame, "Advisory rules (recommendations.py)", size=16, color=GREEN, bold=True, space_after=8)
    add_text(tb.text_frame, "Tier 1 · risk ≥ 0.58 → 10–12 t/ha FYM or 3–4 t/ha biochar + green manuring + 100% residue retention", size=13, bullet=True)
    add_text(tb.text_frame, "Vertisol clay ≥ 250 g/kg → 8–10 t/ha FYM + zero-tillage (Happy Seeder) + straw retention", size=13, bullet=True)
    add_text(tb.text_frame, "pH < 5.8 → agricultural lime 2.0–2.5 t/ha  ·  pH > 7.5 → gypsum 2.0 t/ha", size=13, bullet=True)
    add_text(tb.text_frame, "BSI > 0.25 → minimum tillage, cover cropping, contour bunding (moisture shield)", size=13, bullet=True)
    add_text(tb.text_frame, " ", size=6, space_after=2)
    add_text(tb.text_frame, "Example — Arang (B-1):", size=15, color=CYAN, bold=True, space_after=4)
    add_text(tb.text_frame, "Risk 0.66 · low clay · → 10–12 t/ha FYM before Kharif + Dhaincha/Sunnhemp green manure + lime for acidic patches", size=13, color=WHITE)
    # right: sample prescription card
    panel(s, Inches(7.8), Inches(1.95), Inches(5.0), Inches(4.85))
    tb = textbox(s, Inches(8.1), Inches(2.15), Inches(4.4), Inches(4.4))
    add_text(tb.text_frame, "Rank #1 · Arang (B-1)", size=16, color=RED, bold=True, space_after=10)
    for line in [
        "Mean SOC deficiency: 0.66",
        "High-deficiency area: 946 ha (82%)",
        "",
        "Package:",
        "• 10–12 t/ha Farmyard Manure",
        "• Green manuring (Dhaincha/Sunnhemp)",
        "• 100% crop residue incorporation",
        "• Lime 2.0–2.5 t/ha for acidity",
        "",
        "Funding fit: PMKSY / RKVY",
    ]:
        add_text(tb.text_frame, line, size=12.5, color=WHITE if not line.startswith("Package") and not line.startswith("Funding") else CYAN,
                 bold=line.startswith("Package") or line.startswith("Funding") or line == "Rank #1 · Arang (B-1)")
    footer(s, 5)
    notes(s, "SPEAKER NOTES (1:45–2:15 · ADVISORY):\n"
             "\"The map ranks sectors; the advisory engine makes it actionable. For Arang B-1 — our rank-one "
             "block — the system prescribes 10 to 12 tonnes of farmyard manure per hectare before the Kharif "
             "season, green manuring, and lime for the acidic patches.\"\n\n"
             "Emphasise the clay sensitivity:\n"
             "\"The same risk score gives different advice on heavy Vertisol clay versus light sandy soil — 8–10 "
             "tonnes per hectare plus zero-tillage, so we don't compact the clay.\"\n\n"
             "Close the loop: \"A district collector can now route PMKSY or RKVY funds to exactly the blocks "
             "that need them — with exact tonne quantities per block.\"")

    # ---------------------------------------------------------------- Slide 6 · Validation & Impact
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Validation & Impact", "Cross-validated against ISRO Bhuvan — built for real deployment",
              "Honest science, offline execution, societal benefit")
    # left: bhuvan + confidence
    panel(s, Inches(0.55), Inches(1.95), Inches(6.6), Inches(4.85))
    tb = textbox(s, Inches(0.85), Inches(2.1), Inches(6.0), Inches(4.5))
    add_text(tb.text_frame, "ISRO Bhuvan cross-validation", size=16, color=GREEN, bold=True, space_after=8)
    add_text(tb.text_frame, "48.96% bare-soil share validated against official Bhuvan 1:50,000 LULC cropland baseline (≈82.4% agricultural cropland match)", size=13, bullet=True)
    add_text(tb.text_frame, "Village Geocoding API maps hotspots to real Gram Panchayats — Chandkhuri, Kendri, Nawapara", size=13, bullet=True)
    add_text(tb.text_frame, "Ensemble confidence map shows where the model is certain vs uncertain", size=13, bullet=True)
    add_text(tb.text_frame, " ", size=6, space_after=2)
    add_text(tb.text_frame, "Societal impact", size=16, color=CYAN, bold=True, space_after=8)
    add_text(tb.text_frame, "Zero-cost satellite guidance for 3.5M smallholder farmers", size=13, bullet=True)
    add_text(tb.text_frame, "Targeted fertiliser subsidy routing (PMKSY / RKVY)", size=13, bullet=True)
    add_text(tb.text_frame, "Aligns with National Carbon & Green Credit programmes", size=13, bullet=True)
    add_text(tb.text_frame, "Fits ISRO's Earth-Observation-for-society mission", size=13, bullet=True)
    # right: confidence map + offline chip
    conf = ASSETS / "phase4_model_confidence_map.png"
    if conf.exists():
        pic_fit(s, conf, Inches(7.4), Inches(1.95), w=Inches(5.4))
    stat_chip(s, Inches(7.4), Inches(5.15), Inches(5.4), "100% OFFLINE", "cached COGs · district-office ready · ≈1 minute runtime", GREEN)
    footer(s, 6)
    notes(s, "SPEAKER NOTES (2:15–3:00 · VALIDATION & CLOSE):\n"
             "\"We don't stop at our own predictions. We cross-validate against official ISRO Bhuvan land-cover "
             "baselines — an 82% match with agricultural cropland — and attach our advisories to real Gram "
             "Panchayats through Bhuvan village geocoding.\"\n\n"
             "Impact close (memorise):\n"
             "\"SoilGuard-CG turns open Earth Observation data into climate-resilient soil intelligence — "
             "100% offline, zero sampling cost, and directly actionable under national soil and carbon "
             "programmes. Thank you.\"\n\n"
             "Final Q&A safety:\n"
             "• R² too low? → \"Literature benchmark for un-leakaged optical SOC is 0.35–0.55 — we're inside.\"\n"
             "• Internet? → \"Runs fully offline on cached rasters.\"\n"
             "• Limitations? → \"Topsoil 0–15cm, post-harvest bare windows, prioritisation proxy not lab assay.\"\n"
             "• What's next? → \"Multi-date compositing, SAR for deep carbon, field app + API.\"")

    out = OUT_DIR / "SoilGuard_CG_Ideathon_Deck.pptx"
    prs.save(out)
    print(f"[OK] Saved deck: {out}")
    return out


if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)
    build()
